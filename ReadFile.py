from dotenv import load_dotenv
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import os
import sys
import time
import pdfplumber
from bs4 import BeautifulSoup
import docx
import pandas as pd
from PIL import Image
import fitz  # PyMuPDF for PDF images
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
# Create a vector store with a sample text
# from langchain_core.vectorstores import InMemoryVectorStore
from langchain_ollama import OllamaLLM
#from langchain_core.prompts import PromptTemplate       not required for now
# from langchain.chains import RetrievalQA      why is this used if its not correct what should be used
from pinecone import Pinecone, ServerlessSpec
from langchain_pinecone import PineconeVectorStore

# loading .env file (must be first)
load_dotenv(verbose=True)   # Shows exactly which line fails

# Reconfigure stdout to handle utf-8 characters properly in Windows terminal
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding='utf-8')

file_path = r"C:\Users\saira\Downloads\SAU-I-983-Instructions.pptx"

def ensure_dir(directory):
    """Create directory if it doesn't exist"""
    os.makedirs(directory, exist_ok=True)

def read_file(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    img_dir = f"{base_name}_images"
    ensure_dir(img_dir)
    
    if file_path.lower().endswith('.pdf'):
        # PDF - PyMuPDF for images + pdfplumber for tables/text
        text = ""
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text += f"\n--- PAGE {page_num + 1} ---\n"
            text += page.get_text("text") + "\n"
            
            # Extract images
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                img_path = os.path.join(img_dir, f"pdf_page{page_num+1}_img{img_idx+1}.{image_ext}")
                with open(img_path, 'wb') as img_file:
                    img_file.write(image_bytes)
                text += f"[IMAGE SAVED: {img_path}]\n"
        
        # Tables with pdfplumber (secondary pass)
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables, start=1):
                        if table and len(table) > 1:
                            try:
                                df = pd.DataFrame(table[1:], columns=table[0])
                                text += f"\n[TABLE {page_num}-{table_num}]\n{df.to_string()}\n[/TABLE]\n"
                            except:
                                text += f"[TABLE {page_num}-{table_num} DETECTED]\n"
        
        doc.close()
        return text
        
    elif file_path.lower().endswith('.docx'):
        doc = docx.Document(file_path)
        text = ""
        
        # Text + Tables
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text.strip() + "\n"
        
        for table_idx, table in enumerate(doc.tables, start=1):
            text += f"\n[TABLE {table_idx}]\n"
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text += " | ".join(cells) + "\n"
            text += "[/TABLE]\n"
        
        # EXTRACT IMAGES
        for i, image_part in enumerate(doc.part.blip_store.image_parts):
            image_bytes = image_part.blob
            image_extension = image_part.content_type.split('/')[-1]
            img_path = os.path.join(img_dir, f"doc_img_{i+1}.{image_extension}")
            
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            text += f"[IMAGE SAVED: {img_path}]\n"
        
        return text
        
    elif file_path.lower().endswith('.pptx'):
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            text += f"\n--- SLIDE {slide_num} ---\n"
            
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text += paragraph.text.strip() + "\n"
                elif shape.has_table:
                    text += "[TABLE]\n"
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        text += " | ".join(cells) + "\n"
                    text += "[/TABLE]\n"
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract PPTX images
                    image_obj = shape.image
                    image_bytes = image_obj.blob
                    image_extension = image_obj.content_type.split('/')[-1]
                    img_path = os.path.join(img_dir, f"ppt_slide{slide_num}_img{shape_idx}.{image_extension}")
                    
                    with open(img_path, 'wb') as f:
                        f.write(image_bytes)
                    text += f"[IMAGE SAVED: {img_path}]\n"
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    text += "[CHART DETECTED]\n"
        
        return text
        
    elif file_path.lower().endswith(('.html', '.htm')):
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        soup = BeautifulSoup(html_content, 'html.parser')
        text = soup.get_text()
        
        # Download external images
        for img_num, img in enumerate(soup.find_all('img'), start=1):
            src = img.get('src', '')
            alt = img.get('alt', f'image {img_num}')
            
            if src.startswith('http'):
                try:
                    import requests
                    response = requests.get(src, timeout=5)
                    img_extension = response.headers.get('content-type', '').split('/')[-1] or 'jpg'
                    img_path = os.path.join(img_dir, f"html_img_{img_num}.{img_extension}")
                    
                    with open(img_path, 'wb') as f:
                        f.write(response.content)
                    text += f"[IMAGE SAVED: {img_path}]\n"
                except:
                    text += f"[IMAGE: {alt} | {src} (failed to download)]\n"
            else:
                text += f"[IMAGE: {alt} | {src}]\n"
        return text
        
    else:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except:
            return f"[ERROR: Could not read {os.path.basename(file_path)}]"

def get_text_chunks(text, chunk_size=500, chunk_overlap=50):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    final_split = splitter.split_text(text)
    return final_split

if __name__ == "__main__":
    # Step 1: Read file and split into chunks
    full_text = read_file(file_path)
    print(f"File loaded. Length: {len(full_text)} characters.")
    # print(full_text)
    with open("output.txt", "w", encoding="utf-8") as f:
        f.write(full_text)
    print("extracted text save to output.txt")

    chunks = get_text_chunks(full_text)
    #print(f"Split into {len(chunks)} chunks.")
    for i, chunk in enumerate(chunks):
        if "sevis" in chunk.lower():     # .lower() for case-insensitive search
            print(f"FOUND SEVIS in chunk {i}: {chunk}")
        
    # Step 2: Setup embeddings and LLM
    ##embeddings = OllamaEmbeddings(model="llama3")
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    llm = OllamaLLM(model="llama3")

    # Step 3: Create/connect to Pinecone index (must be done before storing vectors)
    index_name = "my-first-rag-1"
    pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
    if not pc.has_index(index_name):
        pc.create_index(
            name=index_name,
            dimension=384,    # all-MiniLM-L6-v2 dimensions
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1"),
        )

    index = pc.Index(index_name)

    # Step 4: Create vector store and only add text if index is empty
    vectorstore = PineconeVectorStore(index=index, embedding=embeddings)
    
    # Check if we already have vectors in the index to avoid duplicates
    stats = index.describe_index_stats()
    vector_count = stats.get('total_vector_count', 0)

    if vector_count > 0:
        choice = input(f"Index contains {vector_count} vectors. Would you like to clear it and re-index? (y/n): ")
        if choice.lower() == 'y':
            print("Clearing index...")
            index.delete(delete_all=True)
            vector_count = 0
            # Wait a moment for the deletion to propagate
            time.sleep(2)

    if vector_count == 0:
        print(f"Indexing {len(chunks)} chunks...")
        vectorstore.add_texts(texts=chunks)
        print("Indexing complete.")
    else:
        print(f"Skipping indexing. Using existing {vector_count} vectors.")

    # Step 5: Create retriever for querying
    retriever = vectorstore.as_retriever(search_kwargs={"k": 8})
    
    # Step 6: Query loop
    while True:
        question = input("\nEnter your question (or 'quit' to exit): ")

        if question.lower() in ['exit', 'quit', 'bye', 'goodbye', 'q']:
            print("Goodbye!")
            break

        # Retrieve the most similar text
        # retrieved_documents = retriever.invoke(question)
        # Retrieve the most similar text WITH SCORES
        retrieved_documents = vectorstore.similarity_search_with_score(question, k=8)
        print("\n--- SOURCES FOUND ---")
        for i, (doc, score) in enumerate(retrieved_documents, 1):
            print(f"Source {i} (Score: {score:.4f}):\n{doc.page_content[:200]}...")
            
        context = "\n\n".join([doc.page_content for doc, score in retrieved_documents])
        print("\n=== FULL CONTEXT SENT TO LLM ===")
        print(context[:2000])  # Show first 2000 chars
        print("=== END CONTEXT ===\n")

        # Generate response using context with a stricter prompt
        prompt = f"""You are a helpful assistant. Use ONLY the following pieces of context to answer the question at the end. 

        IMPORTANT: Read through ALL the context carefully. The answer might be anywhere in the context, not just at the beginning.
        If you don't know the answer based on the context, just say that you don't know, don't try to make up an answer.

        Context:
        {context}

        Question: {question}

        Think step by step:
        1. What exactly is the question asking for?
        2. Search through the context for the specific answer
        3. Provide only the answer that directly addresses the question
        
        Answer:"""
        
        response = llm.invoke(prompt)
        print(f"\nAnswer: {response}")
else:
    print("Usage: python ReadFile.py")

