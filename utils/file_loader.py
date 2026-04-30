import os
import pdfplumber
import pandas as pd
import fitz  # PyMuPDF
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from datetime import datetime, timezone

def ensure_dir(directory):
    """Create directory if it doesn't exist"""
    os.makedirs(directory, exist_ok=True)

def read_file(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_name = os.path.splitext(os.path.basename(file_path))[0]
    file_type = file_path.split(".")[-1].lower()

    img_dir = f"{file_name}_images"
    ensure_dir(img_dir)
    # This list will hold one dict per page/slide
    # Each dict has: {"text": "...", "metadata": {...}}
    sections = []

    # ─────────────────────────────────────────
    # PDF
    # ─────────────────────────────────────────
    if file_path.lower().endswith(".pdf"):
        # PDF - PyMuPDF for images + pdfplumber for tables/text
        page_text = ""
        doc = fitz.open(file_path)

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text += f"\n--- PAGE {page_num + 1} ---\n"
            page_text += page.get_text("text") + "\n"

            # Extract images
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]

                img_path = os.path.join(
                    img_dir, f"pdf_page{page_num+1}_img{img_idx+1}.{image_ext}"
                )
                with open(img_path, "wb") as img_file:
                    img_file.write(image_bytes)
                page_text += f"[IMAGE SAVED: {img_path}]\n"
            # Save this page as one section with its metadata
            sections.append(
                {
                    "text": page_text,
                    "metadata": {
                        "source": file_name,
                        "file_type": file_type,
                        "page": page_num + 1,  # human-readable (starts at 1)
                        "total_pages": len(doc),
                    },
                }
            )

        # Tables with pdfplumber (secondary pass)
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables, start=1):
                        if table and len(table) > 1:
                            try:
                                df = pd.DataFrame(table[1:], columns=table[0])
                                table_text = df.to_string()
                                sections.append(
                                    {
                                        "text": table_text,
                                        "metadata": {
                                            "source": file_name,
                                            "file_type": file_type,
                                            "page": page_num,
                                            "table_num": table_num,
                                            "content_type": "table",  # extra tag!
                                        },
                                    }
                                )
                            except:
                                pass
        doc.close()

    # ─────────────────────────────────────────
    # DOCX
    # ─────────────────────────────────────────

    elif file_path.lower().endswith(".docx"):
        doc = docx.Document(file_path)

        # All paragraphs together as one section
        para_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        sections.append(
            {
                "text": para_text,
                "metadata": {
                    "source": file_name,
                    "file_type": file_type,
                    "page": 1,  # DOCX has no easy page numbers
                },
            }
        )
        # Each table as its own section
        for table_idx, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(rows)
            sections.append(
                {
                    "text": table_text,
                    "metadata": {
                        "source": file_name,
                        "file_type": file_type,
                        "table_num": table_idx,
                        "content_type": "table",
                    },
                }
            )

        # EXTRACT IMAGES
        try:
            for i, image_part in enumerate(doc.part.blip_store.image_parts):
                img_path = os.path.join(
                    img_dir, f"doc_img_{i+1}.{image_part.content_type.split('/')[-1]}"
                )
                with open(img_path, "wb") as f:
                    f.write(image_part.blob)
        except AttributeError:
            pass  # No images in this document — that is fine

    # ─────────────────────────────────────────
    # PPTX
    # ─────────────────────────────────────────

    elif file_path.lower().endswith(".pptx"):
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = ""

            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text += paragraph.text.strip() + "\n"

                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text += " | ".join(cells) + "\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract PPTX images
                    image_obj = shape.image
                    image_extension = image_obj.content_type.split("/")[-1]
                    img_path = os.path.join(
                        img_dir,
                        f"ppt_slide{slide_num}_img{shape_idx}.{image_extension}",
                    )

                    with open(img_path, "wb") as f:
                        f.write(image_obj.blob)
                    slide_text += f"[IMAGE SAVED: {img_path}]\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    slide_text += "[CHART DETECTED]\n"

            # Save this slide as one section with its metadata
            sections.append(
                {
                    "text": slide_text.strip(),
                    "metadata": {
                        "source": file_name,
                        "file_type": file_type,
                        "slide": slide_num,  # which slide
                        "total_slides": len(prs.slides),
                    },
                }
            )
    # ─────────────────────────────────────────
    # HTML
    # ─────────────────────────────────────────

    elif file_path.lower().endswith((".html", ".htm")):
        with open(file_path, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file.read(), "html.parser")

        # Extract text — preserve table structure
        page_text = ""

        for element in soup.find_all(["p", "h1", "h2", "h3", "h4", "li"]):
            text = element.get_text().strip()
            if text:
                page_text += text + "\n"

        # ── Extract tables with structure preserved ──
        for table in soup.find_all("table"):
            rows = table.find_all("tr")
            for row in rows:
                cells = [cell.get_text().strip() for cell in row.find_all(["td", "th"])]
                if any(cells):  # skip empty rows
                    page_text += " | ".join(cells) + "\n"
            page_text += "\n"

        # Extract images from HTML
        for img_num, img in enumerate(soup.find_all("img"), start=1):
            src = img.get("src", "")
            alt = img.get("alt", f"image {img_num}")

            if src.startswith("http"):
                # Image is on the internet → try to download it
                try:
                    import requests

                    response = requests.get(src, timeout=5)
                    # Get extension from content-type header
                    img_extension = response.headers.get(
                        "content-type", "image/jpg"
                    ).split("/")[-1]
                    img_path = os.path.join(
                        img_dir, f"html_img_{img_num}.{img_extension}"
                    )

                    with open(img_path, "wb") as f:
                        f.write(response.content)
                    page_text += f"[IMAGE SAVED: {img_path}]\n"

                except Exception as e:
                    # If download fails → just note it
                    page_text += f"[IMAGE: {alt} | {src} (failed to download)]\n"
            else:
                # Image is a local file → just note its path
                page_text += f"[IMAGE: {alt} | {src}]\n"

        sections.append(
            {
                "text": page_text,
                "metadata": {
                    "source": file_name,
                    "file_type": file_type,
                    "page": 1,
                },
            }
        )

    else:
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                sections.append(
                    {
                        "text": file.read(),
                        "metadata": {
                            "source": file_name,
                            "file_type": file_type,
                            "page": 1,
                        },
                    }
                )
        except Exception as e:
            sections.append(
                {
                    "text": f"[ERROR reading file: {e}]",
                    "metadata": {"source": file_name, "file_type": file_type},
                }
            )

    # ── Add date_ingested to ALL sections at once ──
    ingestion_date = datetime.now(timezone.utc).strftime(
        "%Y-%m-%d"
    )  # e.g. "2026-04-14"

    for section in sections:
        section["metadata"]["date_ingested"] = ingestion_date

    return sections
