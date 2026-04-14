from dotenv import load_dotenv
from datetime import datetime, timezone, timedelta
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import os
import sys
import time
import json
import re
import pdfplumber
from bs4 import BeautifulSoup
import docx
import pandas as pd
import fitz  # PyMuPDF for PDF images
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_experimental.text_splitter import SemanticChunker
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_ollama import OllamaLLM
from pinecone import Pinecone, ServerlessSpec
from langchain_pinecone import PineconeVectorStore
import cohere

# loading .env file (must be first)
load_dotenv(verbose=True)  # Shows exactly which line fails

# Reconfigure stdout to handle utf-8 characters properly in Windows terminal
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")

# file_path = r"C:\Users\saira\Downloads\SAU-I-983-Instructions.pptx"
# file_path = r"C:\Users\saira\Downloads\sample_onboarding.docx"
file_path = r"C:\Users\saira\Downloads\sample_project_report.html"


def ensure_dir(directory):
    """Create directory if it doesn't exist"""
    os.makedirs(directory, exist_ok=True)


def read_file(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_name = os.path.splitext(os.path.basename(file_path))[0]
    file_type = file_path.split(".")[-1].lower()

    img_dir = f"{file_name}_images"
    ensure_dir(img_dir)
    # This list will hold one dict per page/slide
    # Each dict has: {"text": "...", "metadata": {...}}
    sections = []

    # ─────────────────────────────────────────
    # PDF
    # ─────────────────────────────────────────
    if file_path.lower().endswith(".pdf"):
        # PDF - PyMuPDF for images + pdfplumber for tables/text
        page_text = ""
        doc = fitz.open(file_path)

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text += f"\n--- PAGE {page_num + 1} ---\n"
            page_text += page.get_text("text") + "\n"

            # Extract images
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]

                img_path = os.path.join(
                    img_dir, f"pdf_page{page_num+1}_img{img_idx+1}.{image_ext}"
                )
                with open(img_path, "wb") as img_file:
                    img_file.write(image_bytes)
                page_text += f"[IMAGE SAVED: {img_path}]\n"
            # Save this page as one section with its metadata
            sections.append(
                {
                    "text": page_text,
                    "metadata": {
                        "source": file_name,
                        "file_type": file_type,
                        "page": page_num + 1,  # human-readable (starts at 1)
                        "total_pages": len(doc),
                    },
                }
            )

        # Tables with pdfplumber (secondary pass)
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables, start=1):
                        if table and len(table) > 1:
                            try:
                                df = pd.DataFrame(table[1:], columns=table[0])
                                table_text = df.to_string()
                                sections.append(
                                    {
                                        "text": table_text,
                                        "metadata": {
                                            "source": file_name,
                                            "file_type": file_type,
                                            "page": page_num,
                                            "table_num": table_num,
                                            "content_type": "table",  # extra tag!
                                        },
                                    }
                                )
                            except:
                                pass
        doc.close()

    # ─────────────────────────────────────────
    # DOCX
    # ─────────────────────────────────────────

    elif file_path.lower().endswith(".docx"):
        doc = docx.Document(file_path)

        # All paragraphs together as one section
        para_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        sections.append(
            {
                "text": para_text,
                "metadata": {
                    "source": file_name,
                    "file_type": file_type,
                    "page": 1,  # DOCX has no easy page numbers
                },
            }
        )
        # Each table as its own section
        for table_idx, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(rows)
            sections.append(
                {
                    "text": table_text,
                    "metadata": {
                        "source": file_name,
                        "file_type": file_type,
                        "table_num": table_idx,
                        "content_type": "table",
                    },
                }
            )

        # EXTRACT IMAGES
        try:
            for i, image_part in enumerate(doc.part.blip_store.image_parts):
                img_path = os.path.join(
                    img_dir, f"doc_img_{i+1}.{image_part.content_type.split('/')[-1]}"
                )
                with open(img_path, "wb") as f:
                    f.write(image_part.blob)
        except AttributeError:
            pass  # No images in this document — that is fine

    # ─────────────────────────────────────────
    # PPTX
    # ─────────────────────────────────────────

    elif file_path.lower().endswith(".pptx"):
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = ""

            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text += paragraph.text.strip() + "\n"

                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text += " | ".join(cells) + "\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract PPTX images
                    image_obj = shape.image
                    image_extension = image_obj.content_type.split("/")[-1]
                    img_path = os.path.join(
                        img_dir,
                        f"ppt_slide{slide_num}_img{shape_idx}.{image_extension}",
                    )

                    with open(img_path, "wb") as f:
                        f.write(image_obj.blob)
                    slide_text += f"[IMAGE SAVED: {img_path}]\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    slide_text += "[CHART DETECTED]\n"

            # Save this slide as one section with its metadata
            sections.append(
                {
                    "text": slide_text.strip(),
                    "metadata": {
                        "source": file_name,
                        "file_type": file_type,
                        "slide": slide_num,  # which slide
                        "total_slides": len(prs.slides),
                    },
                }
            )
    # ─────────────────────────────────────────
    # HTML
    # ─────────────────────────────────────────

    elif file_path.lower().endswith((".html", ".htm")):
        with open(file_path, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file.read(), "html.parser")

        # Extract text — preserve table structure
        page_text = ""

        for element in soup.find_all(["p", "h1", "h2", "h3", "h4", "li"]):
            text = element.get_text().strip()
            if text:
                page_text += text + "\n"

        # ── Extract tables with structure preserved ──
        for table in soup.find_all("table"):
            rows = table.find_all("tr")
            for row in rows:
                cells = [cell.get_text().strip() for cell in row.find_all(["td", "th"])]
                if any(cells):  # skip empty rows
                    page_text += " | ".join(cells) + "\n"
            page_text += "\n"

        # Extract images from HTML
        for img_num, img in enumerate(soup.find_all("img"), start=1):
            src = img.get("src", "")
            alt = img.get("alt", f"image {img_num}")

            if src.startswith("http"):
                # Image is on the internet → try to download it
                try:
                    import requests

                    response = requests.get(src, timeout=5)
                    # Get extension from content-type header
                    img_extension = response.headers.get(
                        "content-type", "image/jpg"
                    ).split("/")[-1]
                    img_path = os.path.join(
                        img_dir, f"html_img_{img_num}.{img_extension}"
                    )

                    with open(img_path, "wb") as f:
                        f.write(response.content)
                    page_text += f"[IMAGE SAVED: {img_path}]\n"

                except Exception as e:
                    # If download fails → just note it
                    page_text += f"[IMAGE: {alt} | {src} (failed to download)]\n"
            else:
                # Image is a local file → just note its path
                page_text += f"[IMAGE: {alt} | {src}]\n"

        sections.append(
            {
                "text": page_text,
                "metadata": {
                    "source": file_name,
                    "file_type": file_type,
                    "page": 1,
                },
            }
        )

    else:
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                sections.append(
                    {
                        "text": file.read(),
                        "metadata": {
                            "source": file_name,
                            "file_type": file_type,
                            "page": 1,
                        },
                    }
                )
        except Exception as e:
            sections.append(
                {
                    "text": f"[ERROR reading file: {e}]",
                    "metadata": {"source": file_name, "file_type": file_type},
                }
            )

    # ── Add date_ingested to ALL sections at once ──
    ingestion_date = datetime.now(timezone.utc).strftime(
        "%Y-%m-%d"
    )  # e.g. "2026-04-14"

    for section in sections:
        section["metadata"]["date_ingested"] = ingestion_date

    return sections  # ← returns a LIST of dicts, not a single string:


def get_text_chunks(
    sections, embeddings, file_type=None, chunk_size=1000, chunk_overlap=50
):
    # Choose semantic chunker based on file type
    if file_type in ("pptx", "pdf", "docx"):
        semantic_splitter = SemanticChunker(
            embeddings=embeddings,
            breakpoint_threshold_type="percentile",
            breakpoint_threshold_amount=95,
        )
    elif file_type in ("html", "htm"):
        semantic_splitter = SemanticChunker(
            embeddings=embeddings,
            breakpoint_threshold_type="interquartile",
            breakpoint_threshold_amount=1.5,
        )
    else:
        semantic_splitter = SemanticChunker(
            embeddings=embeddings,
            breakpoint_threshold_type="standard_deviation",
            breakpoint_threshold_amount=3,
        )
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=[
            "\n\n",  # paragraph break → try this first
            "\n",  # single line break → try second
            ". ",  # sentence ending → try third
            " ",  # word boundary → try fourth
            "",  # last resort → cut anywhere
        ],
    )
    all_chunks = []

    for section in sections:
        raw_text = section["text"]
        base_metadata = section["metadata"]

        if not raw_text.strip():
            continue  # skip empty sections

        # Layer 1 → semantic chunking first
        semantic_chunks = semantic_splitter.split_text(raw_text)

        # Layer 2 → split any large chunks further
        for i, sem_chunk in enumerate(semantic_chunks):
            smaller_chunks = splitter.split_text(sem_chunk)

            for j, chunk in enumerate(smaller_chunks):
                all_chunks.append(
                    {
                        "text": chunk,
                        "metadata": {
                            **base_metadata,
                            "chunk_index": f"{i}_{j}",  # e.g. "2_0", "2_1"
                        },
                    }
                )

    return all_chunks


def extract_filters_from_question(question: str) -> dict:
    """
    Uses ollama to automatically detect metadata filters from the user's question.
    Returns a dict like: {"source": "filename", "slide": 3, "date_ingested": "2026-04-14"}
    """

    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    system_prompt = f"""You are a filter extraction assistant. Today's date is {today}.

Your job is to read a user's question and extract search filters from it.

Return ONLY a valid JSON object with these optional keys:
- "source": filename they mention (string, no extension)
- "slide": slide number they mention (integer)
- "content_type": if they mention "table" or "chart" (string)
- "date_ingested": exact date in YYYY-MM-DD format if they mention a specific date (string)
- "days": number of days if they say "last N days" or "past N days" (integer)

Rules:
- Only include a key if the question clearly mentions it
- If nothing matches, return an empty object: {{}}
- Return ONLY the JSON. No explanation. No markdown. No extra text.

Examples:
Question: "What is on slide 3?" → {{"slide": 3}}
Question: "Show me tables from the onboarding file" → {{"source": "sample_onboarding", "content_type": "table"}}
Question: "What did we discuss in the last 7 days?" → {{"days": 7}}
Question: "What is the capital of France?" → {{}}

Now extract filters from this question:
Question: "{question}"
Answer: """

    response = llm.invoke(system_prompt)
    # Clean up response — remove markdown code fences if Ollama adds them
    cleaned = response.strip()
    cleaned = re.sub(r"```json|```", "", cleaned).strip()

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        # If Claude returned something unexpected, return empty (no filters)
        return {}


if __name__ == "__main__":
    # Step 1: Read file and split into chunks
    sections = read_file(file_path)
    file_type = file_path.split(".")[-1].lower()

    # Step 2: Setup embeddings and LLM
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )
    chunks = get_text_chunks(sections, embeddings, file_type)
    # print(chunks)
    llm = OllamaLLM(model="gemma3:4b")

    co = cohere.Client(os.getenv("COHERE_API_KEY"))

    # Step 3: Create/connect to Pinecone index (must be done before storing vectors)
    index_name = "my-first-rag-1"
    pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
    if not pc.has_index(index_name):
        pc.create_index(
            name=index_name,
            # dimension=4096,     #llama3 dimensions (1536 is for OpenAI)
            dimension=384,  # all-MiniLM-L6-v2 dimensions
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1"),
        )

    index = pc.Index(index_name)

    # Step 4: Create vector store and only add text if index is empty
    vectorstore = PineconeVectorStore(index=index, embedding=embeddings)

    # Check if we already have vectors in the index to avoid duplicates
    stats = index.describe_index_stats()
    vector_count = stats.get("total_vector_count", 0)

    if vector_count > 0:
        choice = input(
            f"Index contains {vector_count} vectors. Would you like to clear it and re-index? (y/n): "
        )
        if choice.lower() == "y":
            print("Clearing index...")
            index.delete(delete_all=True)
            vector_count = 0
            # Wait a moment for the deletion to propagate
            time.sleep(2)

    if vector_count == 0:
        print(f"Indexing {len(chunks)} chunks...")
        vectorstore.add_texts(
            texts=[c["text"] for c in chunks], metadatas=[c["metadata"] for c in chunks]
        )
        print("Indexing complete.")
    else:
        print(f"Skipping indexing. Using existing {vector_count} vectors.")

    # Step 5: Create retriever for querying
    retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

    # Step 6: Query loop
    while True:
        question = input("\nEnter your question (or 'quit' to exit): ")

        if question.lower() in ["exit", "quit", "bye", "goodbye", "q"]:
            print("Goodbye!")
            break
        # ── AUTO-DETECT filters from the question ──
        # print(" Detecting filters from your question...")
        detected = extract_filters_from_question(question)

        # if detected:
        #     print(f"  Filters detected: {detected}")
        # else:
        #     print("  No filters detected — searching everything.")

        # ── Build Pinecone filter from detected values ──
        pinecone_filter = {}

        if detected.get("source"):
            pinecone_filter["source"] = {"$eq": detected["source"]}

        if detected.get("slide"):
            pinecone_filter["slide"] = {"$eq": int(detected["slide"])}

        if detected.get("content_type"):
            pinecone_filter["content_type"] = {"$eq": detected["content_type"]}

        if detected.get("date_ingested"):
            pinecone_filter["date_ingested"] = {"$eq": detected["date_ingested"]}

        if detected.get("days"):
            days_ago = int(detected["days"])
            date_range = [
                (datetime.now(timezone.utc) - timedelta(days=i)).strftime("%Y-%m-%d")
                for i in range(days_ago)
            ]
            # Pinecone "$in" means: "match any of these values"
            # Like SQL: WHERE date_ingested IN ('2026-04-14', '2026-04-13', ...)
            pinecone_filter["date_ingested"] = {"$in": date_range}

        # ── If no filters given, search everything ──
        search_filter = pinecone_filter if pinecone_filter else None

        # Add before retrieval
        start = time.time()

        # Retrieve the most similar text WITH SCORES
        retrieved_documents = vectorstore.similarity_search_with_score(
            question, k=20, filter=search_filter
        )

        if not retrieved_documents:
            print("⚠️  No results found with that filter. Try removing the filter.")
            continue

        # # TEMPORARY → just to see the raw output
        # print("\n--- RAW RESULT EXAMPLE ---")
        # doc, score = retrieved_documents[0]  # look at first result only
        # print(f"Text: {doc.page_content[:200]}")
        # print(f"Metadata: {doc.metadata}")
        # print(f"Score: {score}")

        # Re-rank the 10 results using Cohere
        rerank_results = co.rerank(
            query=question,
            documents=[doc.page_content for doc, score in retrieved_documents],
            top_n=5,
            model="rerank-english-v3.0",
        )
        # print(f"rerank_results are {rerank_results}")

        # Build reranked list in new order
        reranked_documents = [
            retrieved_documents[result.index] for result in rerank_results.results
        ]

        retrieval_time = time.time()

        # print("\n--- SOURCES FOUND ---")
        # for i, (doc, score) in enumerate(retrieved_documents, 1):
        #     meta = doc.metadata
        #     location = meta.get("slide") or meta.get("page") or "?"
        #     print(
        #         f"  [{i}] Score: {score:.3f} | "
        #         f"File: {meta.get('source')} | "
        #         f"Page/Slide: {location} | "
        #         f"Chunk: {meta.get('chunk_index')}"
        #     )

        context = "\n\n".join([doc.page_content for doc, score in reranked_documents])

        print(f"\n📊 Context stats:")
        # print(context)
        print(f"  Chunks: {len(retrieved_documents)}")
        print(f"  Characters: {len(context):,}")
        print(f"  Estimated tokens: ~{len(context) // 4:,}")

        # Generate response using context with a stricter prompt
        prompt = f"""You are a helpful assistant. Use ONLY the following pieces of context to answer the question at the end.

        IMPORTANT: Read through ALL the context carefully. The answer might be anywhere in the context, not just at the beginning.
        If you don't know the answer based on the context, just say that you don't know, don't try to make up an answer.

        Context:
        {context}

        Question: {question}

        Think step by step:
        1. What exactly is the question asking for?
        2. Search through the context for the specific answer
        3. Provide only the answer that directly addresses the question

        Answer:"""

        response = llm.invoke(prompt)
        # generation_time = time.time()
        # print(f"\n⏱️ TIMING:")
        # print(f"  Retrieval: {retrieval_time - start:.2f}s")
        # print(f"  Generation: {generation_time - retrieval_time:.2f}s")
        # print(f"  Total: {generation_time - start:.2f}s")
        print(f"\nAnswer: {response}")
else:
    print("Usage: python ReadFile.py")
